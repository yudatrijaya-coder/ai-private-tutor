import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625) # 16:9 widescreen ratio

# Colors
NAVY = RGBColor(15, 23, 42)        # Slate 900
ICE_BLUE = RGBColor(186, 230, 253) # Sky 200
ACCENT_BLUE = RGBColor(14, 165, 233) # Sky 500
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(51, 65, 85)   # Slate 700
LIGHT_BG = RGBColor(248, 250, 252) # Slate 50
CONTAINER_BG = RGBColor(255, 255, 255)
BORDER_COLOR = RGBColor(226, 232, 240) # Slate 200

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def format_text_frame(tf):
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

# ==============================================================================
# SLIDE 1: Title Slide (Dark Theme)
# ==============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1, NAVY)

# Background visual accent (large thin rectangle border)
border_rect = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(9.2), Inches(4.825))
border_rect.fill.background()
border_rect.line.color.rgb = RGBColor(51, 65, 85)
border_rect.line.width = Pt(1.5)

# Main Title & Subtitle in one text frame
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(2.2))
tf1 = title_box.text_frame
format_text_frame(tf1)

p_title = tf1.paragraphs[0]
p_title.text = "AI PRIVATE TUTOR"
p_title.font.name = "Georgia"
p_title.font.size = Pt(46)
p_title.font.bold = True
p_title.font.color.rgb = WHITE
p_title.space_after = Pt(8)

p_sub = tf1.add_paragraph()
p_sub.text = "Personalized Telegram-Integrated LMS & Student Helper Platform"
p_sub.font.name = "Calibri"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = ICE_BLUE
p_sub.space_after = Pt(20)

# URL & Bot Info Box
info_box = slide1.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(8.4), Inches(1.2))
tf_info = info_box.text_frame
format_text_frame(tf_info)

p_url = tf_info.paragraphs[0]
p_url.text = "🌐 Web Portal: https://senangbelajar.web.id"
p_url.font.name = "Consolas"
p_url.font.size = Pt(15)
p_url.font.color.rgb = ACCENT_BLUE
p_url.space_after = Pt(6)

p_bot = tf_info.add_paragraph()
p_bot.text = "🤖 Telegram Bot: @senangbelajar_bot"
p_bot.font.name = "Consolas"
p_bot.font.size = Pt(15)
p_bot.font.color.rgb = ACCENT_BLUE

# ==============================================================================
# SLIDE 2: Tech Stack & Architecture (Light Theme with Diagram)
# ==============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2, LIGHT_BG)

# Slide Title
title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.6))
tf2 = title2_box.text_frame
format_text_frame(tf2)
p2_title = tf2.paragraphs[0]
p2_title.text = "Platform Architecture & Tech Stack"
p2_title.font.name = "Georgia"
p2_title.font.size = Pt(26)
p2_title.font.bold = True
p2_title.font.color.rgb = NAVY

# Left Column: Key Stack Details
left_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(3.8))
tf_left = left_box.text_frame
format_text_frame(tf_left)

stack_items = [
    ("Next.js 14 (App Router)", "Powers the main responsive student dashboard & admin controls."),
    ("Telegraf.js Webhooks", "Handles instant dual-channel interactions via @senangbelajar_bot."),
    ("Prisma & PostgreSQL", "Provides reliable, type-safe database schemas with migration tracking."),
    ("SumoPod LLM Gateway", "Runs deepseek-v4-flash for cost-efficient AI private tutoring.")
]

for idx, (tech, desc) in enumerate(stack_items):
    p = tf_left.paragraphs[0] if idx == 0 else tf_left.add_paragraph()
    p.space_after = Pt(12)
    
    run_tech = p.add_run()
    run_tech.text = f"✔  {tech}\n"
    run_tech.font.name = "Calibri"
    run_tech.font.size = Pt(14)
    run_tech.font.bold = True
    run_tech.font.color.rgb = NAVY

    run_desc = p.add_run()
    run_desc.text = desc
    run_desc.font.name = "Calibri"
    run_desc.font.size = Pt(13)
    run_desc.font.color.rgb = DARK_TEXT

# Right Column: Visual Diagram
diagram_left = 5.2
diagram_top = 1.3
diagram_width = 4.3

# Box 1: Telegram User
tg_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(diagram_left), Inches(diagram_top), Inches(diagram_width), Inches(0.7))
tg_box.fill.solid()
tg_box.fill.fore_color.rgb = NAVY
tg_box.line.color.rgb = ACCENT_BLUE
tg_box.text = "Telegram Client (Student Chat)"
tg_box.text_frame.paragraphs[0].font.name = "Calibri"
tg_box.text_frame.paragraphs[0].font.size = Pt(13)
tg_box.text_frame.paragraphs[0].font.bold = True

# Box 2: Next.js Webhook & DB
app_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(diagram_left), Inches(diagram_top + 1.2), Inches(diagram_width), Inches(1.1))
app_box.fill.solid()
app_box.fill.fore_color.rgb = CONTAINER_BG
app_box.line.color.rgb = BORDER_COLOR
tf_app = app_box.text_frame
tf_app.word_wrap = True
p_app1 = tf_app.paragraphs[0]
p_app1.text = "Next.js backend (Webhook /api/bot)"
p_app1.font.bold = True
p_app1.font.size = Pt(12)
p_app1.font.color.rgb = NAVY
p_app2 = tf_app.add_paragraph()
p_app2.text = "DB: PostgreSQL | ORM: Prisma"
p_app2.font.size = Pt(11)
p_app2.font.color.rgb = DARK_TEXT

# Box 3: SumoPod LLM
llm_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(diagram_left), Inches(diagram_top + 2.8), Inches(diagram_width), Inches(0.7))
llm_box.fill.solid()
llm_box.fill.fore_color.rgb = CONTAINER_BG
llm_box.line.color.rgb = ACCENT_BLUE
llm_box.text = "SumoPod LLM Gateway (deepseek-v4-flash)"
llm_box.text_frame.paragraphs[0].font.name = "Calibri"
llm_box.text_frame.paragraphs[0].font.size = Pt(12)
llm_box.text_frame.paragraphs[0].font.bold = True
llm_box.text_frame.paragraphs[0].font.color.rgb = NAVY

# ==============================================================================
# SLIDE 3: Core Features (2x2 Grid)
# ==============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3, LIGHT_BG)

# Title
title3_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.6))
tf3 = title3_box.text_frame
format_text_frame(tf3)
p3_title = tf3.paragraphs[0]
p3_title.text = "Core Platform Features"
p3_title.font.name = "Georgia"
p3_title.font.size = Pt(26)
p3_title.font.bold = True
p3_title.font.color.rgb = NAVY

features = [
    ("Dual Curriculum Integration", "Unifies Kemendikdasmen's SIBI (national textbooks) and school-specific Moodle portals into a single learning path.", 0.5, 1.2),
    ("YouTube Video Linker", "Scrapes, filters, and maps high-quality reference videos to curriculum topics. Validated via oEmbed API.", 5.1, 1.2),
    ("Adaptive Quiz Engine", "Generates and reviews student knowledge using dynamic question banks generated directly matching progress.", 0.5, 3.2),
    ("React Flow Mindmaps", "Automatically visualizes curriculum outlines and study guides into interactive node maps on the student dashboard.", 5.1, 3.2)
]

for title, desc, col_x, row_y in features:
    # Card shape
    card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(col_x), Inches(row_y), Inches(4.4), Inches(1.6))
    card.fill.solid()
    card.fill.fore_color.rgb = CONTAINER_BG
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.5)
    
    # Left accent colored bar
    accent_bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(col_x), Inches(row_y), Inches(0.08), Inches(1.6))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar.line.fill.background()

    # Content
    content_box = slide3.shapes.add_textbox(Inches(col_x + 0.2), Inches(row_y + 0.1), Inches(4.1), Inches(1.4))
    tf_c = content_box.text_frame
    format_text_frame(tf_c)
    
    p_head = tf_c.paragraphs[0]
    p_head.text = title
    p_head.font.name = "Calibri"
    p_head.font.size = Pt(15)
    p_head.font.bold = True
    p_head.font.color.rgb = NAVY
    p_head.space_after = Pt(4)
    
    p_desc = tf_c.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Calibri"
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = DARK_TEXT

# ==============================================================================
# SLIDE 4: Operations & Default Schedules (Dark Theme)
# ==============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4, NAVY)

# Slide Title
title4_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.6))
tf4 = title4_box.text_frame
format_text_frame(tf4)
p4_title = tf4.paragraphs[0]
p4_title.text = "Operations & Default Schedule"
p4_title.font.name = "Georgia"
p4_title.font.size = Pt(26)
p4_title.font.bold = True
p4_title.font.color.rgb = WHITE

# Left Panel: Schedule (using container shape)
sched_container = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(4.3), Inches(3.6))
sched_container.fill.solid()
sched_container.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
sched_container.line.color.rgb = RGBColor(71, 85, 105) # Slate 600
tf_sched = sched_container.text_frame
format_text_frame(tf_sched)

p_sc_title = tf_sched.paragraphs[0]
p_sc_title.text = "📅 Default Study Schedule"
p_sc_title.font.bold = True
p_sc_title.font.size = Pt(16)
p_sc_title.font.color.rgb = ICE_BLUE
p_sc_title.space_after = Pt(14)

schedules = [
    ("Intensive Schedule", "Senin, Rabu, Jumat\n19:00 - 21:00 WIB (UTC+7)"),
    ("Reguler Schedule", "Selasa, Kamis, Minggu\n19:00 - 20:00 WIB (UTC+7)"),
    ("Notification Policy", "Reminder sent exactly 15 mins before (18:45 WIB) with Topic and Sub-topic selection.")
]

for idx, (title, details) in enumerate(schedules):
    p = tf_sched.add_paragraph()
    p.space_after = Pt(10)
    
    run_t = p.add_run()
    run_t.text = f"• {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(13)
    run_t.font.color.rgb = WHITE
    
    run_d = p.add_run()
    run_d.text = details
    run_d.font.size = Pt(12)
    run_d.font.color.rgb = ICE_BLUE

# Right Panel: Operational Automation
ops_box = slide4.shapes.add_textbox(Inches(5.1), Inches(1.2), Inches(4.4), Inches(3.6))
tf_ops = ops_box.text_frame
format_text_frame(tf_ops)

ops_items = [
    ("Moodle Automation Pipeline", "Automated scrapers fetch Moodle PDFs, metadata, and lesson structures securely using user enrolment tokens."),
    ("Discrepancy Checking Engine", "Periodic background jobs check physical directories against database tables to match student-specific courses."),
    ("Robust Hosting Infrastructure", "Managed PM2 application runners and Caddy server with SSL handling incoming API calls and bot webhooks.")
]

for idx, (title, details) in enumerate(ops_items):
    p = tf_ops.paragraphs[0] if idx == 0 else tf_ops.add_paragraph()
    p.space_after = Pt(14)
    
    run_t = p.add_run()
    run_t.text = f"✔  {title}\n"
    run_t.font.bold = True
    run_t.font.size = Pt(14)
    run_t.font.color.rgb = ICE_BLUE
    
    run_d = p.add_run()
    run_d.text = details
    run_d.font.size = Pt(12)
    run_d.font.color.rgb = RGBColor(226, 232, 240)

# Footer URL line across all slides
footer_box = slide4.shapes.add_textbox(Inches(0.5), Inches(5.15), Inches(9.0), Inches(0.4))
tf_foot = footer_box.text_frame
p_foot = tf_foot.paragraphs[0]
p_foot.alignment = PP_ALIGN.RIGHT
p_foot.text = "https://senangbelajar.web.id  |  @senangbelajar_bot"
p_foot.font.name = "Consolas"
p_foot.font.size = Pt(10)
p_foot.font.color.rgb = RGBColor(148, 163, 184)

# Save
prs.save("ai_private_tutor_project_resume.pptx")
print("Presentation updated successfully with premium design and structural layout!")
